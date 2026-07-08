# ============================================
# 备课助手 Web 版 - 用户登录版
# 启动: python web_app.py
# 访问: http://localhost:5002
# ============================================

import re
import os
import sys
import yaml
import random
import sqlite3
import threading
import queue
from pathlib import Path
from flask import Flask, render_template, request, jsonify, send_file, Response, redirect, url_for, make_response
from openai import OpenAI
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
import time
import json

# 导入数据库模块
import database

def get_base_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys._MEIPASS)
    return Path(__file__).parent

def get_data_dir():
    if getattr(sys, 'frozen', False):
        return Path(sys.executable).parent
    return Path(__file__).parent

BASE_DIR = get_base_dir()
DATA_DIR = get_data_dir()

template_dir = str(BASE_DIR / "templates")
app = Flask(__name__, template_folder=template_dir)
app.secret_key = 'beike_assistant_secret_key_2024'
app.config['TEMPLATES_AUTO_RELOAD'] = True

# 读配置
CONFIG_PATH = BASE_DIR / "config.yaml"
cfg = {}
if CONFIG_PATH.exists():
    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

API_KEY = cfg.get("api_key", "sk-xxx")
BASE_URL = cfg.get("base_url", "https://api.deepseek.com/v1")
MODEL = cfg.get("model", "deepseek-chat")
TEMP = cfg.get("temperature", 0.7)
MAX_TOKENS = cfg.get("max_tokens", 3072)
MAX_HISTORY = cfg.get("max_history", 6)
OUTPUT_DIR = cfg.get("output_dir", ".")
if not os.path.isabs(OUTPUT_DIR):
    OUTPUT_DIR = os.path.abspath(os.path.join(str(DATA_DIR), OUTPUT_DIR))
else:
    OUTPUT_DIR = os.path.abspath(OUTPUT_DIR)

os.makedirs(OUTPUT_DIR, exist_ok=True)

HISTORY_DIR = os.path.join(OUTPUT_DIR, "历史记录")
os.makedirs(HISTORY_DIR, exist_ok=True)

def get_user_output_dir(user_id=None):
    if user_id:
        user_dir = os.path.join(HISTORY_DIR, f"user_{user_id}")
    else:
        user_dir = os.path.join(HISTORY_DIR, "public")
    os.makedirs(user_dir, exist_ok=True)
    return user_dir

SYSTEM_PROMPT = """你是专业的「AI 备课助手」，服务对象是中小学 / 职业院校的一线教师。你的任务是根据教师的需求，按规范格式生成高质量的教学资源，并**严格用指定标签包裹可保存的文件内容**，方便系统自动解析并下载。

---

## 🎯 核心能力（6 大模块）
1. **PPT 课件生成**：输出结构化 PPT 大纲，系统会自动解析为 `.pptx` 课件
2. **PPT 讲解纲要**：为每页幻灯片生成详细的教师讲解稿（逐页讲稿 + 互动设计）
3. **教案生成**：45 分钟标准格式教案（三维目标 / 教学过程 / 板书 / 反思）
4. **习题生成**：基础 / 提高 / 拓展 三级难度分层习题，附答案解析
5. **作业分析**：作业批改数据分析 + 错因归因 + 辅导建议 + 教学调整策略
6. **复习提纲**：期末 / 单元复习资料，知识脉络 + 核心考点 + 典型例题

---

## 🚫 最重要的输出规则（不遵守会导致文件无法下载）
**必须用 `[文件:xxx.md]  …内容…  [/文件]` 或 `[PPT:xxx.pptx]  …大纲…  [/PPT]` 标签把生成结果包裹起来**，系统会自动保存为可下载的文件。
- 如果用户要的是「教案 / 习题 / 作业分析 / 讲解纲要 / 复习提纲」 → 用 `[文件:xxx.md]`
- 如果用户要的是「PPT 课件 / 幻灯片」 → 用 `[PPT:xxx.pptx]`
- 若用户同时要多份内容（如「同时生成教案和PPT」），可同时出现多个标签块
- 纯文字的简短说明 / 寒暄可以写在标签外，但**核心生成内容必须全部放在标签块内**

✅ 正确示例：
```
好的，为您生成《背影》的教案如下：
[文件:初中语文《背影》教案.md]
# 初中语文《背影》教案

## 一、三维教学目标
- **知识与技能**：……
- **过程与方法**：……
- **情感态度与价值观**：……

## 二、教学重难点
……（以下内容全部写在标签内）……

## 六、教学反思栏
……
[/文件]
```

❌ 错误示例（严禁！会导致用户无法下载文件）：
- 直接输出 Markdown，不写任何标签包裹
- 标签名写错，如 `[教案:xxx]`、`[下载:xxx]`
- 文件内容一半在标签里，一半在标签外

---

## 📖 模块 1：教案生成 标准格式（`[文件:xxx教案.md]`）
结构必须完整，标题层级严格按下面顺序：

```
# 《课题名》教案（学段 + 学科）

## 一、基本信息
- 学科：xxx
- 学段：xxx（小学/初中/高中/中职）
- 年级：xxx
- 课时：1课时（45分钟）
- 课型：新授课 / 复习课 / 讲评课 / 实验课 / 活动课

## 二、三维教学目标 / 核心素养目标
- 知识与技能：……
- 过程与方法：……
- 情感态度与价值观：……

## 三、教学重难点
- **教学重点**：……
- **教学难点**：……

## 四、学情分析与教学方法
- 学情分析：……
- 教学方法：讲授法 / 讨论法 / 实验法 / 情境教学法 / 任务驱动法

## 五、教学准备
- 教具准备：多媒体课件、实验器材、学案……
- 学生准备：……

## 六、详细教学过程（精确到分钟）
### 1. 导入新课（约 5 分钟）
……教师活动……
……学生活动……

### 2. 新知讲授（约 15 分钟）
……分步骤讲解……

### 3. 课堂练习 / 小组讨论（约 12 分钟）
……

### 4. 课堂小结（约 5 分钟）
……由学生总结 + 教师补充……

### 5. 作业布置（约 3 分钟）
- 基础作业：……
- 拓展作业：……

### 6. 板书设计（约 5 分钟同步板书）
主板书 | 副板书
-------|-------
……    | ……

## 七、板书设计
- 主板书：……（结构化呈现）
- 副板书：……

## 八、分层作业设计
- ✅ **基础层（必做）**：……
- 📈 **提高层（选做）**：……
- 🚀 **拓展层（挑战）**：……

## 九、教学反思
- 本节课亮点：……
- 待改进点：……
- 课堂生成问题记录（留白教师填写）：
```

---

## 📝 模块 2：习题生成 标准格式（`[文件:xxx习题.md]`）
```
# 《课题 / 知识点》分层练习题

## 一、基本信息
- 学段学科：xxx
- 考查知识点：xxx
- 题目总数：xx 道
- 建议用时：xx 分钟
- 难度分布：基础 60% / 提高 30% / 拓展 10%

---

## 二、基础巩固题（约 60%）
### 一、选择题（每题 3 分，共 xx 分）
1. 题目内容……
   A. xxx   B. xxx   C. xxx   D. xxx
2. ……

### 二、填空题（每题 2 分，共 xx 分）
1. _________________
2. _________________

---

## 三、能力提高题（约 30%）
### 三、解答题 / 计算题（每题 xx 分，共 xx 分）
1. 题目内容……

---

## 四、拓展探究题（约 10%）
### 四、综合应用题 / 探究题
1. 题目内容……

---

## 五、参考答案与详细解析
1. 选择题 1：答案 B
   - 解析：……
   - 考点：……
   - 易错分析：……
2. ……
```

---

## 📊 模块 3：作业分析报告（`[文件:xxx作业分析.md]`）
```
# 《作业名称》学情分析报告

## 一、整体概况
- 学科：xxx
- 班级：xxx（共 xx 人）
- 应交：xx 份，实交：xx 份，上交率：xx%
- 平均分：xx 分 ｜ 最高分：xx ｜ 最低分：xx ｜ 及格率：xx% ｜ 优秀率：xx%

### 分数段分布
| 分数段    | 人数 | 占比  |
|-----------|------|-------|
| 90-100 分 | x 人 | xx%  |
| 80-89 分  | x 人 | xx%  |
| 70-79 分  | x 人 | xx%  |
| 60-69 分  | x 人 | xx%  |
| 60 分以下 | x 人 | xx%  |

## 二、知识点掌握情况雷达
| 知识点         | 平均得分率 | 掌握等级 |
|----------------|------------|----------|
| 知识点 A       | xx%        | 熟练 / 一般 / 薄弱 |
| 知识点 B       | xx%        | ……     |

## 三、高频错题 TOP 5 + 归因分析
| 排名 | 题号 | 知识点 | 得分率 | 主要错解 | 错因归类 |
|------|------|--------|--------|----------|----------|
| 1    | 第x题 | xxx | xx% | …… | ⚠️ 概念不清 / 审题错误 / 计算失误 / 方法未掌握 |
| 2    | …… | …… | …… | …… | …… |

### 典型错解展示（第 x 题）
> 错误解法示例：……
> ✅ 正确解法：……
> 🔍 错因剖析：……

## 四、分层辅导建议
### 🔴 学困生（xx 人，60 分以下）
- 薄弱点：……
- 辅导策略：① …… ② …… ③ ……
- 补充练习题：……

### 🟡 中等生（xx 人，60-84 分）
- 提升点：……
- 辅导策略：……
- 强化练习：……

### 🟢 优等生（xx 人，85 分以上）
- 拓展方向：……
- 挑战性题目：……

## 五、下节课教学调整建议
- 需要补讲的知识点：……（建议用时 xx 分钟）
- 课堂讲评顺序建议：先讲第 x/x/x 题（得分率最低）
- 教学方法调整：……
- 与家长沟通要点（家校共育）：……

## 六、讲评课时分配建议
| 环节            | 时间  | 内容说明 |
|-----------------|-------|----------|
| 整体情况通报    | 3 分钟 | 成绩分布 + 表扬优秀进步学生 |
| 高频错题精讲    | 20 分钟 | TOP 5 错因分析 + 变式训练 |
| 小组互助订正    | 10 分钟 | 学困生结对，同伴讲解 |
| 针对性补充练习  | 8 分钟  | 相似题型，当堂巩固 |
| 总结 + 二次过关 | 4 分钟  | 关键方法总结，小测过关 |
```

---

## 🎤 模块 4：PPT 讲解纲要（逐页讲稿，`[文件:xxx讲解纲要.md]`）
```
# 《课件名》PPT 逐页讲解纲要

## 使用说明
- 总页数：xx 页
- 建议总时长：45 分钟
- 适用对象：xxx 年级学生

---

## 第 1 页：封面（约 1 分钟）
🎤 **教师讲解词**：
同学们好，今天我们一起来学习《……》。在上课之前请大家看屏幕上的这张图片 / 这个问题，有没有同学能说一说……
❓ **互动提问**（可选）：……
⭐ **语气提示**：语速稍慢，吸引注意力

---

## 第 2 页：教学目标（约 2 分钟）
🎤 **教师讲解词**：
本节课我们要达成三个学习目标：第一，……；第二，……；第三，……
⭐ **强调**：第 2 条是重点，请用红色笔在学案上画出来
⏰ 本页用时：2 分钟

---

## 第 3 页：xxx（按 PPT 每页依次写）
🎤 **教师讲解词**：……
❓ **互动提问**：……（预留学生作答的留白，可写学生可能的回答）
⚠️ **易错提醒**：……
⏰ 本页用时：xx 分钟

---
（每页 PPT 都按上面格式写一节）
---

## 最后一页：课堂小结 + 作业布置（约 3 分钟）
🎤 **教师讲解词**：
好，本节课我们学习了……，主要内容可以用三句话记住：①……②……③……
下课后请大家完成：……（作业内容）
下课，同学们再见！
```

---

## 🤖 模块 5：复习提纲（`[文件:xxx复习提纲.md]`）
```
# 《课程名》期末 / 单元复习提纲

## 📚 一、复习内容概览
- 章节范围：第 x 章 — 第 x 章
- 建议复习用时：xx 课时
- 重要程度：★★★★★（5 星为必考）

## 🧠 二、核心知识体系（思维导图式）
### 模块一：xxx
- 核心概念 1：……
  - 定义：……
  - 关键词：……
- 核心概念 2：……

### 模块二：xxx
……

## 📐 三、重点公式 / 定理 / 结论速记
| 名称 | 公式 / 结论 | 适用条件 | 考频 |
|------|-------------|----------|------|
| xxx  | ……         | ……       | ⭐⭐⭐ |

## 📝 四、典型例题精讲
### 例题 1（★★ 基础题）
> 题目：……
> 解题步骤：
> ① ……
> ② ……
> 💡 点拨：……

### 例题 2（★★★★ 高频考点）
……

## ⚠️ 五、易错点警示 TOP 10
1. ❌ 错误理解：…… → ✅ 正确：……
2. ❌ ……

## 🎯 六、考点预测 + 分值分布
| 考点         | 预测分值 | 题型       | 难度 |
|--------------|----------|------------|------|
| xxx          | 8-12 分  | 选择 + 解答 | ★★★ |

## ✅ 七、自我检测（附答案）
……（10 道精选小题，覆盖全部考点）
```

---

## 📊 模块 6：PPT 课件生成（必须用 `[PPT:xxx.pptx]` … `[/PPT]` 包裹）
**结构严格遵守以下格式，系统会自动解析成 `.pptx` 文件：**

```
[PPT:荷塘月色_语文课件.pptx]
# 《荷塘月色》语文课件
朱自清散文 · 高中语文必修上册

## 第1页：情境导入
- 展示荷花池夜色图片
- 提问：同学们记忆中关于"月"和"荷"的诗句有哪些？
- 引出作者：朱自清
- 板书课题

## 第2页：学习目标
- 📖 知识与技能：掌握重点词语，理解关键语句含义
- 🧠 过程与方法：通过朗读体会文章的语言美和意境美
- ❤️ 情感态度：体会作者情感，感受文学作品的感染力

## 第3页：作者介绍
- 朱自清（1898-1948），原名自华，字佩弦
- 现代著名散文家、诗人、学者
- 代表作品：《背影》《春》《匆匆》《荷塘月色》
- 散文风格：语言洗练，文笔清丽，情感真挚

## 第4页：写作背景
……

## 第N页：课堂小结 + 作业布置
……
[/PPT]
```

**PPT 格式要求（务必遵守）：**
1. 第一行 `# 主标题` 是封面主标题，第二行（无特殊符号）是副标题
2. 每页幻灯片用 `## 第N页：页面标题` 开头
3. 每页下面的内容用 `- 要点` 列（最多 7 条，每条简短）
4. 页数控制在 8-20 页之间，结构完整
5. **绝对不要**把多页内容塞进一个 `##` 里

---

## 🎨 其它输出规范
- 全文中文表达，专业简洁
- Markdown 层级清晰：# 一级 → ## 二级 → ### 三级
- 重点用 **粗体**、表格整理数据、列表梳理条目
- 编号必须连续，避免乱码字符

## 📎 上传文件智能识别规则（最重要！）
用户消息中如果出现 `【用户上传文件：文件名（类型）】` 标记，说明用户上传了该文件的全文内容，你必须按以下逻辑处理：

### 情况 A：文件类型 = "教案/教学设计模板"（文件名或内容含「教案」「教学设计」「学案」「模板」「____」「[填空]」「__」等占位符）
- 你的任务 = **按模板中所有的空白 / 占位符 / 待填写项，自动填充完整的教学内容**
- 规则：
  1. 严格**保留模板原有的章节结构、编号、表格结构、填空占位符的位置**
  2. 模板里写着「____」「___」「...」「待填写」「填写」「留白」「XXX」「[     ]」的位置，全部填上具体、详实、符合学科逻辑的内容
  3. 若用户在上传文件之外还写了需求（如"请按人教版高二语文"），需优先按该需求填充内容
  4. 没有的学科/年级信息时，按模板里能推断的学段和学科默认选一个最通用的
  5. **生成结果必须用 `[文件:xxx_已填充.md]` 完整包裹起来**，把整个模板+填充后的全部内容放进去，方便用户下载

### 情况 B：文件类型 = "学生作业/成绩/批改记录"（文件名或内容含「作业」「成绩」「得分」「分数」「错题」「批改」「考勤」「班级」「学生」「姓名」等）
- 你的任务 = **按「模块 5 作业分析」标准格式生成一份完整的学情分析报告**
- 报告内容必须包含 6 部分：
  1. **整体概况**（班级人数、均分、及格率、优秀率、最高分/最低分、分数段分布）
  2. **错题归因分析**（哪题错最多？错因归类：概念不清 / 计算失误 / 审题错误 / 知识点没掌握 / 书写规范）
  3. **知识点掌握雷达**（按知识点列出 掌握率 = 做对人数/总人数）
  4. **学生分层画像**（优生/中等/待进步各占多少 %，每一层的典型问题）
  5. **针对性辅导建议**（全班教学调整 + 分层布置作业 + 个别辅导名单）
  6. **后续教学调整策略**（下一节课怎么改、是否加小测、是否安排讲评课）
- **生成结果必须用 `[文件:xxx_作业分析报告.md]` 完整包裹**

### 情况 C：其他类型文件（如参考资料、课文原文、讲义等）
- 先明确告诉用户："已读取到您上传的《xxx》（共 x 字 / x 行）"
- 再结合用户的具体问题给出对应的回答；如果用户没有明确问题，自动判断：内容像课文原文 → 建议生成《xxx》教案/PPT/习题；内容像复习资料 → 建议生成复习提纲

最后再次强调：**所有生成内容必须放在对应的标签块中，否则用户点击"下载"时将没有任何文件！**"""

# 存会话（内存缓存 + 数据库持久化）
sessions: dict[str, list[dict]] = {}

def get_current_user():
    token = request.cookies.get('session_token')
    if token:
        return database.get_user_by_token(token)
    return None

def trim_history(messages):
    if not messages or len(messages) <= MAX_HISTORY + 1:
        return messages
    system_msg = messages[0]
    recent = messages[-(MAX_HISTORY):]
    return [system_msg] + recent

def markdown_to_word(markdown_text: str, title: str = "教案") -> BytesIO:
    doc = Document()
    style = doc.styles['Normal']
    font = style.font
    font.name = '微软雅黑'
    font.size = Pt(11)

    title_para = doc.add_paragraph()
    title_run = title_para.add_run(title)
    title_run.font.size = Pt(18)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(31, 78, 121)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    lines = markdown_text.split('\n')
    in_code_block = False
    code_content = []

    for line in lines:
        if line.strip().startswith('```'):
            if in_code_block:
                if code_content:
                    code_para = doc.add_paragraph()
                    code_run = code_para.add_run('\n'.join(code_content))
                    code_run.font.name = 'Courier New'
                    code_run.font.size = Pt(9)
                    code_run.font.color.rgb = RGBColor(128, 128, 128)
                    code_para.paragraph_format.left_indent = Pt(20)
                code_content = []
                in_code_block = False
            else:
                in_code_block = True
            continue

        if in_code_block:
            code_content.append(line)
            continue

        if line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('#### '):
            doc.add_heading(line[5:], level=4)
        elif '**' in line:
            parts = line.split('**')
            para = doc.add_paragraph()
            for i, part in enumerate(parts):
                if i % 2 == 1:
                    run = para.add_run(part)
                    run.bold = True
                else:
                    para.add_run(part)
        elif line.strip().startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.strip().startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif re.match(r'^\d+\.\s', line.strip()):
            doc.add_paragraph(re.sub(r'^\d+\.\s', '', line.strip()), style='List Number')
        elif '|' in line and line.strip():
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if len(cells) >= 2:
                if not all(re.match(r'^-+$', cell) for cell in cells):
                    para = doc.add_paragraph()
                    para.add_run(line)
        elif line.strip().startswith('>'):
            para = doc.add_paragraph(line[1:].strip())
            para.paragraph_format.left_indent = Pt(20)
            para.runs[0].font.color.rgb = RGBColor(128, 128, 128)
        elif line.strip() == '---':
            doc.add_paragraph('_' * 50)
        elif line.strip():
            doc.add_paragraph(line)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def _strip_md(text):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'`(.+?)`', r'\1', text)
    text = re.sub(r'^#{1,6}\s*', '', text)
    text = text.replace('•', '').replace('◦', '').replace('▪', '')
    text = re.sub(r'^\s*[-*]\s+', '', text)
    text = re.sub(r'^\s*\d+[.、)]\s*', '', text)
    return text.strip()


def _detect_header_level(s):
    m = re.match(r'^(#{1,6})\s+(.+)$', s)
    if m:
        return len(m.group(1)), m.group(2).strip()
    return 0, s


def _is_bullet(s):
    return bool(re.match(r'^\s*[-*•◦▪]\s+', s))


def _is_numbered(s):
    return bool(re.match(r'^\s*\d+[.、)]\s*', s))


def _indent_level(raw):
    stripped = raw.lstrip(' \t')
    n = len(raw) - len(stripped)
    if n == 0:
        return 0
    if n <= 3:
        return 1
    if n <= 6:
        return 2
    return 3


def parse_ppt_outline(markdown_text):
    lines = markdown_text.split('\n')
    title_main = ""
    title_sub = ""
    raw_sections = []

    header_levels = set()
    for raw in lines:
        s = raw.strip()
        if not s:
            continue
        lvl, _ = _detect_header_level(s)
        if lvl > 0:
            header_levels.add(lvl)

    if header_levels:
        if 1 in header_levels:
            cover_level = 1
            slide_levels = {l for l in header_levels if l > 1}
        else:
            cover_level = 0
            slide_levels = header_levels
    else:
        cover_level = 0
        slide_levels = set()

    min_slide = min(slide_levels) if slide_levels else 0

    cur_title = None
    cur_items = []

    def flush():
        nonlocal cur_title, cur_items
        if cur_title is not None:
            raw_sections.append((cur_title, cur_items))
        cur_title = None
        cur_items = []

    for raw in lines:
        if not raw.strip():
            continue
        s = raw.strip()
        lvl, content = _detect_header_level(s)

        if lvl == cover_level and cover_level > 0:
            flush()
            title_main = content
        elif lvl in slide_levels:
            flush()
            t = re.sub(r'^第\s*\d+\s*页\s*[：:]\s*', '', content)
            cur_title = t
        elif lvl > 0 and min_slide > 0:
            rel = max(0, lvl - min_slide)
            cur_items.append((min(rel, 3), content))
        elif _is_bullet(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*[-*•◦▪]\s+', '', s)
            cur_items.append((ind, content))
        elif _is_numbered(s):
            ind = _indent_level(raw)
            content = re.sub(r'^\s*\d+[.、)]\s*', '', s)
            cur_items.append((ind, content))
        else:
            if title_main and not title_sub and not raw_sections and cur_title is None:
                title_sub = s
            else:
                cur_items.append((_indent_level(raw), s))
    flush()

    merged = []
    for stitle, items in raw_sections:
        if not items and merged:
            prev_t, prev_items = merged[-1]
            prev_items.append((0, stitle))
            merged[-1] = (prev_t, prev_items)
        else:
            merged.append((stitle, items))
    raw_sections = merged

    if not raw_sections and title_main:
        all_items = []
        for raw in lines:
            s = raw.strip()
            if not s or _detect_header_level(s)[0] == cover_level:
                continue
            if _is_bullet(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*[-*•◦▪]\s+', '', s)))
            elif _is_numbered(s):
                all_items.append((_indent_level(raw), re.sub(r'^\s*\d+[.、)]\s*', '', s)))
            elif _detect_header_level(s)[0] > 0:
                _, content = _detect_header_level(s)
                all_items.append((0, content))
            else:
                all_items.append((_indent_level(raw), s))
        if all_items:
            raw_sections = [("内容概览", all_items)]

    return title_main, title_sub, raw_sections


def _split_overflow(items, max_items=7):
    if len(items) <= max_items:
        return [items]
    result = []
    cur = []
    for lvl, text in items:
        if len(cur) >= max_items and lvl == 0:
            result.append(cur)
            cur = []
        cur.append((lvl, text))
    if cur:
        result.append(cur)
    return result


class PptTheme:
    PRIMARY = PptRGBColor(0x1E, 0x40, 0xAF)
    DARK = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT = PptRGBColor(0x1F, 0x29, 0x37)
    TEXT_SEC = PptRGBColor(0x4B, 0x55, 0x63)
    MUTED = PptRGBColor(0x9C, 0xA3, 0xAF)
    WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = PptRGBColor(0x3B, 0x82, 0xF6)
    FONT = 'Microsoft YaHei'


class PptBuilder:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.sw = self.prs.slide_width
        self.sh = self.prs.slide_height
        self.m = Inches(0.8)

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _tbox(self, slide, left, top, w, h, anchor=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tf

    def _run(self, para, text, size, color, bold=False):
        run = para.add_run()
        run.text = text
        run.font.size = PptPt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = PptTheme.FONT
        return run

    def _rect(self, slide, left, top, w, h, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def _page_num(self, slide, num, total):
        tf = self._tbox(slide, self.sw - Inches(1.6), self.sh - Inches(0.45),
                        Inches(1.4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        self._run(p, f"{num} / {total}", 11, PptTheme.MUTED)

    def _footer(self, slide):
        tf = self._tbox(slide, self.m, self.sh - Inches(0.45),
                        Inches(4), Inches(0.3), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        self._run(p, "AI 备课助手", 10, PptTheme.MUTED)

    def _font_size(self, items):
        n = len(items)
        chars = sum(len(_strip_md(t)) for _, t in items)
        max_lvl = max((lvl for lvl, _ in items), default=0)
        if n <= 4 and chars <= 120:
            base = 30
        elif n <= 6 and chars <= 250:
            base = 26
        elif n <= 8 and chars <= 400:
            base = 22
        elif n <= 12 and chars <= 600:
            base = 18
        else:
            base = 15
        if max_lvl >= 2 and base > 22:
            base = 22
        return base

    def cover(self, title, subtitle, speaker='', date_str=''):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(2.8), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.1), self.sw - self.m * 2,
                        Inches(2), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        t = _strip_md(title) or "课件"
        size = 46 if len(t) <= 18 else (36 if len(t) <= 36 else 28)
        self._run(p, t, size, PptTheme.WHITE, bold=True)

        if subtitle:
            tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                            Inches(0.8), MSO_ANCHOR.TOP)
            p = tf.paragraphs[0]
            self._run(p, _strip_md(subtitle), 22, PptTheme.MUTED)

        footer = '  |  '.join(filter(None, [speaker, date_str]))
        if footer:
            tf = self._tbox(slide, self.m, self.sh - Inches(0.8),
                            self.sw - self.m * 2, Inches(0.4), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            self._run(p, footer, 14, PptTheme.MUTED)

    def content(self, title, items, page_num, total, dark=False):
        slide = self._blank()

        bg = PptTheme.DARK if dark else PptTheme.WHITE
        tc = PptTheme.WHITE if dark else PptTheme.TEXT
        sc = PptTheme.MUTED if dark else PptTheme.TEXT_SEC
        bc = PptTheme.ACCENT if dark else PptTheme.PRIMARY

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg

        tf = self._tbox(slide, self.m, Inches(0.5), self.sw - self.m * 2,
                        Inches(1), MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        clean_t = _strip_md(title)
        ts = 36 if len(clean_t) <= 20 else (28 if len(clean_t) <= 40 else 22)
        self._run(p, clean_t, ts, tc, bold=True)

        self._rect(slide, self.m, Inches(1.5), Inches(2), Inches(0.06), bc)

        base = self._font_size(items)

        BULLETS = ["●", "○", "▪", "·"]
        INDENT_EMU = [0, 457200, 914400, 1371600]
        SIZE_DEC = [0, 3, 5, 6]

        valid_items = [(min(lvl, 3), _strip_md(text)) for lvl, text in items
                       if _strip_md(text)]
        n = len(valid_items)
        est_line_h = (base * 1.3 + 8) / 72.0
        est_height = Inches(min(n * est_line_h + 0.4, 4.6))

        avail_top = Inches(1.85)
        avail_bottom = self.sh - Inches(0.7)
        avail_h = avail_bottom - avail_top
        if est_height < avail_h:
            content_top = avail_top + (avail_h - est_height) / 2
            content_h = est_height
            anchor = MSO_ANCHOR.MIDDLE
        else:
            content_top = avail_top
            content_h = avail_h
            anchor = MSO_ANCHOR.TOP

        tf = self._tbox(slide, self.m, content_top, self.sw - self.m * 2,
                        content_h, anchor)

        first = True
        for lvl, clean in valid_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = PptPt(6)
            p.line_spacing = 1.3

            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(INDENT_EMU[lvl]))
            pPr.set('indent', str(-228600))

            bchar = BULLETS[lvl]
            bsize = base - SIZE_DEC[lvl]
            tsize = base - SIZE_DEC[lvl]
            tcolor = tc if lvl == 0 else sc
            self._run(p, f"{bchar}  ", bsize, bc, bold=(lvl == 0))
            self._run(p, clean, tsize, tcolor, bold=False)

        self._page_num(slide, page_num, total)
        if not dark:
            self._footer(slide)

    def ending(self):
        slide = self._blank()
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptTheme.DARK

        self._rect(slide, self.m, Inches(3.0), Inches(1.2), Inches(0.1), PptTheme.PRIMARY)

        tf = self._tbox(slide, self.m, Inches(3.3), self.sw - self.m * 2,
                        Inches(1.5), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "谢谢观看", 48, PptTheme.WHITE, bold=True)

        tf = self._tbox(slide, self.m, Inches(4.8), self.sw - self.m * 2,
                        Inches(0.6), MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._run(p, "AI 备课助手 · 自动生成", 16, PptTheme.MUTED)

    def build(self, title, subtitle, content_slides, speaker='', date_str=''):
        self.cover(title, subtitle, speaker, date_str)

        expanded = []
        for stitle, items in content_slides:
            chunks = _split_overflow(items, max_items=7)
            for ci, chunk in enumerate(chunks):
                ct = stitle if ci == 0 else f"{stitle}（续{ci}）"
                expanded.append((ct, chunk))

        total = len(expanded) + 1
        for i, (stitle, items) in enumerate(expanded):
            is_summary = any(k in stitle for k in
                             ['总结', '小结', '结语', '谢谢', '感谢', 'CTA', '行动'])
            self.content(stitle, items, i + 2, total, dark=is_summary)

        if not expanded:
            self.content("内容概览", [(0, "暂无具体内容")], 2, 2)

        self.ending()

        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf


def markdown_to_ppt(markdown_text, title="课件"):
    title_main, title_sub, slides = parse_ppt_outline(markdown_text)
    if not title_main:
        title_main = title
    if not slides:
        slides = [(title_main or title,
                   [(0, t) for t in markdown_text.split('\n') if t.strip()])]

    cover_data = {}
    content_slides = []
    for stitle, items in slides:
        if stitle == '封面' or stitle.startswith('封面'):
            cover_data = _parse_kv(items)
        else:
            content_slides.append((stitle, items))

    if cover_data:
        c_title = _strip_md(cover_data.get('主标题', title_main))
        c_sub = _strip_md(cover_data.get('副标题', title_sub))
        speaker = _strip_md(cover_data.get('演讲人', ''))
        date_str = _strip_md(cover_data.get('日期', time.strftime('%Y年%m月%d日')))
    else:
        c_title = _strip_md(title_main) if title_main else title
        c_sub = _strip_md(title_sub) if title_sub else ''
        speaker = ''
        date_str = time.strftime('%Y年%m月%d日')

    builder = PptBuilder()
    return builder.build(c_title, c_sub, content_slides, speaker, date_str)


def generate_ppt_files(reply_text, user_id=None):
    base_dir = get_user_output_dir(user_id)
    ppt_blocks = re.findall(r"\[PPT:\s*([^\]]+)\](.*?)\[/PPT\]", reply_text, re.DOTALL)
    saved_ppts = []
    for filename, outline in ppt_blocks:
        filename = filename.strip()
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"
        safe_name = re.sub(r'[<>:"/\\|?*]', '', filename)
        try:
            title = os.path.splitext(safe_name)[0]
            ppt_buffer = markdown_to_ppt(outline, title=title)
            filepath = os.path.join(base_dir, safe_name)
            with open(filepath, "wb") as f:
                f.write(ppt_buffer.getvalue())
            saved_ppts.append(safe_name)
        except Exception as e:
            print(f"[PPT] 生成失败 {safe_name}: {e}")
    display = reply_text
    if ppt_blocks:
        display = re.sub(r"\[PPT:\s*[^\]]+\]", "", display)
        display = re.sub(r"\[/PPT\]", "", display)
    return saved_ppts, display


def auto_wrap_and_save_fallback(user_message: str, ai_reply: str, user_id):
    """兜底：AI 忘记写 [文件:] / [PPT:] 标签时，自动检测内容类型并保存文件，保证一定能下载。
    返回 (extra_files, extra_ppts, display)"""
    import os as _os, re as _re, time as _time
    if not ai_reply or not ai_reply.strip():
        return [], [], ai_reply
    has_file = bool(_re.search(r"\[文件:\s*[^\]]+\]", ai_reply))
    has_ppt = bool(_re.search(r"\[PPT:\s*[^\]]+\]", ai_reply))
    if has_file or has_ppt:
        return [], [], ai_reply

    text = (user_message or "") + " " + (ai_reply or "")
    lower = text.lower()

    safe_title = _re.sub(r'[<>:"/\\|?*]', '', (user_message or "AI生成内容").strip())[:20] or "AI生成内容"
    ts = _time.strftime("%Y%m%d_%H%M%S")

    if ("ppt" in lower) or ("课件" in text) or ("幻灯片" in text):
        filename = f"课件_{safe_title}_{ts}.pptx"
        wrapped = f"[PPT:{filename}]\n{ai_reply.strip()}\n[/PPT]"
        saved_ppts, display = generate_ppt_files(wrapped, user_id)
        return [], saved_ppts, display

    if "作业" in text and ("分析" in text or "批改" in text or "错因" in text or "讲评" in text):
        fn = f"作业分析_{safe_title}_{ts}.md"
    elif "教案" in text or "教学设计" in text or "说课稿" in text or "教学过程" in text:
        fn = f"教案_{safe_title}_{ts}.md"
    elif "习题" in text or "练习" in text or "试卷" in text or "测试题" in text or ("题" in text and "答案" in text):
        fn = f"习题_{safe_title}_{ts}.md"
    elif ("讲解" in text and ("纲要" in text or "讲稿" in text or "逐页" in text)) or "说课稿" in text:
        fn = f"讲解纲要_{safe_title}_{ts}.md"
    elif "复习" in text or "提纲" in text or "知识点总结" in text or "知识体系" in text:
        fn = f"复习提纲_{safe_title}_{ts}.md"
    else:
        fn = f"备课资料_{safe_title}_{ts}.md"

    base_dir = get_user_output_dir(user_id)
    filepath = _os.path.join(base_dir, fn)
    content = ai_reply.strip()
    if content.startswith("```"):
        content = _re.sub(r"^```\w*\s*\n", "", content)
        content = _re.sub(r"\n```\s*$", "", content)
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    return [fn], [], ai_reply


def extract_and_save_all(user_message, reply_text, user_id):
    """统一处理 AI 回复：解析文件块 + 生成 PPT + 兜底保存，返回(saved_files, saved_ppts, display_for_chat)"""
    user_output_dir = get_user_output_dir(user_id)

    saved_files = []
    files = re.findall(r"\[文件:\s*([^\]]+)\](.*?)\[/文件\]", reply_text, re.DOTALL)
    for filename, content in files:
        filepath = os.path.join(user_output_dir, filename)
        clean = content.strip()
        if clean.startswith("```"):
            clean = re.sub(r"^```\w*\s*\n", "", clean)
            clean = re.sub(r"\n```\s*$", "", clean)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(clean)
        saved_files.append(filename)

    display = reply_text
    if files:
        display = re.sub(r"\[文件:.*?\[/文件\]", "", reply_text, flags=re.DOTALL).strip()

    saved_ppts, display = generate_ppt_files(display, user_id)

    if not saved_files and not saved_ppts:
        extra_fs, extra_ppts, display = auto_wrap_and_save_fallback(
            user_message, reply_text, user_id
        )
        saved_files.extend(extra_fs)
        saved_ppts.extend(extra_ppts)

    return saved_files, saved_ppts, display


@app.route("/")
def index():
    user = get_current_user()
    if not user:
        return redirect(url_for('login'))
    return render_template("index.html", model=MODEL, username=user['username'])

@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        data = request.json if request.is_json else request.form
        username = data.get("username")
        password = data.get("password")
        
        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400
        
        user = database.authenticate_user(username, password)
        if user:
            token = database.generate_session_token()
            database.set_session_token(user['id'], token)
            response = jsonify({"status": "success", "username": user['username']})
            response.set_cookie('session_token', token, httponly=True, secure=False, max_age=86400*7)
            return response
        else:
            return jsonify({"error": "用户名或密码错误"}), 401
    
    return render_template("login.html")

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        data = request.json if request.is_json else request.form
        username = data.get("username")
        password = data.get("password")
        email = data.get("email", "")
        
        if not username or not password:
            return jsonify({"error": "用户名和密码不能为空"}), 400
        
        if len(password) < 6:
            return jsonify({"error": "密码长度至少6位"}), 400
        
        success = database.register_user(username, password, email)
        if success:
            user = database.authenticate_user(username, password)
            if user:
                token = database.generate_session_token()
                database.set_session_token(user['id'], token)
                response = jsonify({"status": "success", "username": user['username']})
                response.set_cookie('session_token', token, httponly=True, secure=False, max_age=86400*7)
                return response
        else:
            return jsonify({"error": "用户名已存在"}), 409
    
    return render_template("register.html")

@app.route("/logout")
def logout():
    token = request.cookies.get('session_token')
    if token:
        database.clear_session_token(token)
    response = make_response(redirect(url_for('login')))
    response.set_cookie('session_token', '', expires=0)
    return response

@app.route("/api/chat", methods=["POST"])
def chat():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = request.json
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    session_id = str(data.get("session_id", "default"))[:64]
    user_message = data.get("message", "").strip()
    attachments = data.get("attachments") or []

    if isinstance(attachments, list):
        blocks = []
        for att in attachments:
            if not isinstance(att, dict):
                continue
            name = str(att.get("name", "未命名文件"))[:120]
            kind = str(att.get("kind", "资料"))[:30]
            text = str(att.get("text", ""))[:MAX_TEXT_CHARS]
            blocks.append(
                f"【用户上传文件：{name}（{kind}）】\n"
                f"文件内容预览（节选，共约 {len(text)} 字）：\n"
                f"```\n{text}\n```\n"
                f"【上传文件结束】"
            )
        if blocks:
            if user_message:
                user_message = user_message + "\n\n---\n\n" + "\n\n".join(blocks)
            else:
                user_message = "\n\n".join(blocks)

    if not user_message:
        return jsonify({"error": "Message cannot be empty"}), 400
    
    if len(user_message) > 24000:
        user_message = user_message[:24000] + "\n\n[内容过长已截断]"

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return jsonify({"error": "Invalid temperature or max_tokens value"}), 400

    if session_id not in sessions:
        sessions[session_id] = [{"role": "system", "content": SYSTEM_PROMPT}]

    sessions[session_id].append({"role": "user", "content": user_message})

    database.save_search_history(user['id'], user_message)

    try:
        client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
        msgs = trim_history(sessions[session_id])
        response = client.chat.completions.create(
            model=_model,
            messages=msgs,
            temperature=_temp,
            max_tokens=_max_tokens,
        )
    except Exception as e:
        return jsonify({"error": f"API request failed: {str(e)}"}), 500

    reply = response.choices[0].message.content
    usage = response.usage

    sessions[session_id].append({"role": "assistant", "content": reply})

    first_user_msg = next((m["content"] for m in sessions[session_id] if m["role"] == "user"), "")
    title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
    database.save_user_chat_session(user['id'], session_id, title)

    saved_files, saved_ppts, display = extract_and_save_all(user_message, reply, user['id'])

    return jsonify({
        "reply": display,
        "raw_reply": reply,
        "model": response.model,
        "usage": {
            "prompt_tokens": usage.prompt_tokens,
            "completion_tokens": usage.completion_tokens,
        },
        "saved_files": saved_files,
        "saved_ppts": saved_ppts,
    })

def extract_keywords(message):
    keywords = []
    subjects = ["物理", "数学", "语文", "英语", "化学", "生物", "历史", "地理", "政治", "信息技术", "通用技术", "体育", "音乐", "美术", "科学", "道德与法治", "思想品德"]
    for subj in subjects:
        if subj in message:
            keywords.append(subj)
            break
    levels = ["小学", "初中", "高中", "大学", "职业院校", "中职", "高职", "一年级", "二年级", "三年级", "四年级", "五年级", "六年级", "初一", "初二", "初三", "高一", "高二", "高三"]
    for level in levels:
        if level in message:
            keywords.append(level)
            break
    types = ["教案", "课件", "习题", "试卷", "复习", "提纲", "实验", "实训", "说课稿", "导学案", "任务单", "教学设计"]
    for t in types:
        if t in message:
            keywords.append(t)
            break
    if not keywords:
        keywords = ["教学", "备课"]
    return keywords

knowledge_base = {
    "物理": ["牛顿运动定律", "能量守恒", "电磁感应", "光学", "热力学", "力学", "相对论", "量子物理", "波动", "磁场"],
    "数学": ["函数", "方程", "几何", "概率统计", "数列", "三角函数", "向量", "导数", "积分", "不等式"],
    "语文": ["文言文", "现代文", "诗词鉴赏", "写作", "阅读", "修辞手法", "表现手法", "作文", "记叙文", "议论文"],
    "英语": ["语法", "词汇", "阅读", "写作", "听力", "口语", "时态", "句型", "翻译", "完形填空"],
    "化学": ["元素周期表", "化学反应", "有机化学", "无机化学", "化学平衡", "电化学", "化学实验", "化学键"],
    "生物": ["细胞", "遗传", "生态", "光合作用", "呼吸作用", "进化论", "微生物", "人体生理"],
    "历史": ["中国古代史", "中国近代史", "世界史", "历史事件", "历史人物", "朝代", "战争", "改革"],
    "地理": ["自然地理", "人文地理", "气候", "地形", "地图", "环境保护", "区域地理", "人口城市"],
    "政治": ["经济生活", "政治生活", "文化生活", "哲学生活", "法律", "道德", "国情"],
    "default": ["教学目标", "教学重难点", "教学过程", "教学方法", "板书设计", "课后作业", "教学反思", "学情分析"]
}

def generate_thinking_phases(user_message):
    keywords = extract_keywords(user_message)
    subject = keywords[0] if keywords else "default"
    kbs = knowledge_base.get(subject, knowledge_base["default"])
    selected_kb = random.sample(kbs, min(4, len(kbs)))

    phases = [
        {"title": "📋 需求分析阶段", "items": [
            f"解析用户需求：{user_message[:25]}{'...' if len(user_message) > 25 else ''}",
            f"识别学科与学段：{'、'.join(keywords)}",
            "明确输出类型与格式要求",
        ]},
        {"title": "🔍 知识检索阶段", "items": [
            "在知识库中检索相关教学资源...",
            f"匹配知识点：{selected_kb[0]}、{selected_kb[1]}",
            "查找课程标准与教学大纲要求",
            f"筛选参考资料：{selected_kb[2]}相关案例",
        ]},
        {"title": "🧠 内容规划阶段", "items": [
            "梳理知识体系与逻辑结构",
            "设计三维教学目标",
            "确定教学重点与难点",
            "规划教学环节与时间分配",
            "选择教学方法与策略",
        ]},
        {"title": "✍️ 内容生成阶段", "items": [
            "组织语言，生成详细内容...",
            "优化表述，确保专业准确",
            "调整结构，保证逻辑清晰",
            "检查内容完整性与合理性",
        ]},
    ]
    return phases

def thinking_to_text(phases, phase_idx, item_idx):
    lines = []
    for i, phase in enumerate(phases):
        if i > phase_idx:
            break
        lines.append(f"【{phase['title']}】")
        for j, item in enumerate(phase['items']):
            if i == phase_idx and j > item_idx:
                break
            if i < phase_idx or (i == phase_idx and j < item_idx):
                icon = "✓"
            elif i == phase_idx and j == item_idx:
                icon = "⏳"
            else:
                icon = "○"
            lines.append(f"  {icon} {item}")
        lines.append("")
    return "\n".join(lines)

@app.route("/api/chat/stream", methods=["POST"])
def chat_stream():
    user = get_current_user()
    if not user:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': '请先登录'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    data = request.json
    if not data or not isinstance(data, dict):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Invalid request body'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    session_id = str(data.get("session_id", "default"))[:64]
    user_message = data.get("message", "").strip()
    attachments = data.get("attachments") or []
    
    if isinstance(attachments, list):
        blocks = []
        for att in attachments:
            if not isinstance(att, dict):
                continue
            name = str(att.get("name", "未命名文件"))[:120]
            kind = str(att.get("kind", "资料"))[:30]
            text = str(att.get("text", ""))[:MAX_TEXT_CHARS]
            blocks.append(
                f"【用户上传文件：{name}（{kind}）】\n"
                f"文件内容预览（节选，共约 {len(text)} 字）：\n"
                f"```\n{text}\n```\n"
                f"【上传文件结束】"
            )
        if blocks:
            if user_message:
                user_message = user_message + "\n\n---\n\n" + "\n\n".join(blocks)
            else:
                user_message = "\n\n".join(blocks)

    if not user_message:
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Message cannot be empty'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )
    
    if len(user_message) > 24000:
        user_message = user_message[:24000] + "\n\n[内容过长已截断]"

    _model = data.get("model", MODEL)
    try:
        _temp = max(0.0, min(2.0, float(data.get("temperature", TEMP))))
        _max_tokens = max(1, min(16384, int(data.get("max_tokens", MAX_TOKENS))))
    except (ValueError, TypeError):
        return Response(
            f"data: {json.dumps({'type': 'error', 'content': 'Invalid temperature or max_tokens value'}, ensure_ascii=False)}\n\n",
            mimetype="text/event-stream"
        )

    if session_id not in sessions:
        sessions[session_id] = [{"role": "system", "content": SYSTEM_PROMPT}]

    sessions[session_id].append({"role": "user", "content": user_message})
    database.save_search_history(user['id'], user_message)

    def generate():
        thinking_phases = generate_thinking_phases(user_message)
        thinking_done = False
        ai_done = False
        ai_error = None
        full_reply = ""
        ai_queue = queue.Queue()
        first_content_received = False

        def ai_worker():
            nonlocal full_reply, ai_done, ai_error
            try:
                client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
                msgs = trim_history(sessions[session_id])
                stream = client.chat.completions.create(
                    model=_model,
                    messages=msgs,
                    temperature=_temp,
                    max_tokens=_max_tokens,
                    stream=True,
                )
                for chunk in stream:
                    if chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        full_reply += content
                        ai_queue.put(("content", content))
                ai_queue.put(("done", None))
                ai_done = True
            except Exception as e:
                ai_error = str(e)
                ai_queue.put(("error", str(e)))
                ai_done = True

        ai_thread = threading.Thread(target=ai_worker)
        ai_thread.start()

        total_thinking_time = 0
        min_thinking_time = 1.2
        max_thinking_time = 3.0

        for phase_idx, phase in enumerate(thinking_phases):
            if thinking_done:
                break
            for item_idx, item in enumerate(phase["items"]):
                if thinking_done:
                    break

                thinking_text = thinking_to_text(thinking_phases, phase_idx, item_idx)
                yield f"data: {json.dumps({'type': 'thinking', 'content': thinking_text}, ensure_ascii=False)}\n\n"

                sleep_time = random.uniform(0.25, 0.4)
                if total_thinking_time + sleep_time > max_thinking_time and phase_idx >= 2:
                    sleep_time = 0.1
                total_thinking_time += sleep_time
                time.sleep(sleep_time)

                if total_thinking_time >= min_thinking_time:
                    while not ai_queue.empty():
                        msg_type, msg_data = ai_queue.get()
                        if msg_type == "content":
                            first_content_received = True
                            break
                        elif msg_type == "done":
                            first_content_received = True
                            ai_done = True
                            break
                        elif msg_type == "error":
                            first_content_received = True
                            ai_done = True
                            break

                if first_content_received and total_thinking_time >= min_thinking_time:
                    thinking_done = True
                    final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
                    yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
                    break

        if not thinking_done:
            final_thinking = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
            yield f"data: {json.dumps({'type': 'thinking_done', 'content': final_thinking}, ensure_ascii=False)}\n\n"
            thinking_done = True

        if ai_error:
            yield f"data: {json.dumps({'type': 'error', 'content': ai_error}, ensure_ascii=False)}\n\n"
            return

        while not ai_done or not ai_queue.empty():
            try:
                msg_type, msg_data = ai_queue.get(timeout=0.1)
                if msg_type == "content":
                    yield f"data: {json.dumps({'type': 'content', 'content': msg_data}, ensure_ascii=False)}\n\n"
                elif msg_type == "done":
                    ai_done = True
            except queue.Empty:
                if not ai_done:
                    time.sleep(0.05)

        sessions[session_id].append({"role": "assistant", "content": full_reply})

        first_user_msg = next((m["content"] for m in sessions[session_id] if m["role"] == "user"), "")
        title = first_user_msg[:30] + "..." if len(first_user_msg) > 30 else first_user_msg
        database.save_user_chat_session(user['id'], session_id, title)

        saved_files, saved_ppts, display = extract_and_save_all(user_message, full_reply, user['id'])

        final_thinking_text = thinking_to_text(thinking_phases, len(thinking_phases) - 1, len(thinking_phases[-1]["items"]) - 1)
        yield f"data: {json.dumps({'type': 'done', 'content': display, 'raw_content': full_reply, 'saved_files': saved_files, 'saved_ppts': saved_ppts, 'thinking': final_thinking_text}, ensure_ascii=False)}\n\n"

    return Response(generate(), mimetype="text/event-stream")

@app.route("/api/clear", methods=["POST"])
def clear():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = request.json
    session_id = data.get("session_id", "default")
    sessions[session_id] = [{"role": "system", "content": SYSTEM_PROMPT}]
    return jsonify({"status": "ok"})

@app.route("/api/sessions", methods=["GET"])
def list_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    db_sessions = database.get_user_sessions(user['id'])
    result = []
    for s in db_sessions:
        result.append({
            "id": s['session_id'],
            "title": s['title'],
            "message_count": 0,
            "created_at": s['created_at']
        })
    return jsonify({"sessions": result})

@app.route("/api/session/<session_id>", methods=["GET"])
def get_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    if session_id in sessions:
        msgs = sessions[session_id]
        user_msgs = [m for m in msgs if m["role"] == "user"]
        title = "新对话"
        if user_msgs:
            first_user = user_msgs[0]["content"]
            title = first_user[:20] + "..." if len(first_user) > 20 else first_user
        return jsonify({
            "id": session_id,
            "title": title,
            "messages": msgs
        })
    
    db_sessions = database.get_user_sessions(user['id'])
    db_session = next((s for s in db_sessions if s['session_id'] == session_id), None)
    if db_session:
        return jsonify({
            "id": session_id,
            "title": db_session['title'],
            "messages": []
        })
    
    return jsonify({"error": "会话不存在"}), 404

@app.route("/api/session/<session_id>", methods=["DELETE"])
def delete_session(session_id):
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    database.delete_user_session(user['id'], session_id)
    if session_id in sessions:
        del sessions[session_id]
    return jsonify({"status": "ok"})

@app.route("/api/sessions/batch", methods=["DELETE"])
def delete_sessions_batch():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = request.json
    ids = data.get("ids", [])
    deleted = 0
    for sid in ids:
        database.delete_user_session(user['id'], sid)
        if sid in sessions:
            del sessions[sid]
        deleted += 1
    return jsonify({"status": "ok", "deleted": deleted})

@app.route("/api/sessions/all", methods=["DELETE"])
def delete_all_sessions():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    sessions.clear()
    return jsonify({"status": "ok"})

@app.route("/api/search_history", methods=["GET"])
def get_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    history = database.get_search_history(user['id'], limit=50)
    return jsonify({"history": history})

@app.route("/api/search_history", methods=["DELETE"])
def clear_search_history():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute('DELETE FROM search_history WHERE user_id = ?', (user['id'],))
    conn.commit()
    conn.close()
    return jsonify({"status": "ok"})

@app.route("/api/config", methods=["GET"])
def get_config():
    return jsonify({
        "model": MODEL,
        "base_url": BASE_URL,
        "temperature": TEMP,
        "max_tokens": MAX_TOKENS,
        "output_dir": OUTPUT_DIR,
    })

MIME_MAP = {
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.md': 'text/markdown; charset=utf-8',
    '.txt': 'text/plain; charset=utf-8',
    '.csv': 'text/csv; charset=utf-8',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    '.pdf': 'application/pdf',
}

ALLOWED_UPLOAD_EXT = {'.md', '.txt', '.docx', '.xlsx', '.csv', '.pdf'}
MAX_UPLOAD_BYTES = 20 * 1024 * 1024  # 20MB
MAX_TEXT_CHARS = 30000  # 单文件最多送3万字符给AI


def _read_text_bytes(b: bytes) -> str:
    for enc in ('utf-8-sig', 'utf-8', 'gbk', 'gb18030', 'utf-16'):
        try:
            return b.decode(enc)
        except UnicodeDecodeError:
            continue
    return b.decode('utf-8', errors='ignore')


def _guess_file_kind(name: str, text: str) -> str:
    low = (name + "\n" + text[:1500]).lower()
    template_keys = ['教案', '教学设计', '学案', '模板', '____', '___', '待填写', '留白', '[    ]', '[___]', '填写说明', '教学目标', '教学重难点']
    homework_keys = ['作业', '成绩', '得分', '分数', '错题', '批改', '班级', '学生', '姓名', '学号', '排名', '及格', '均分', '考勤', '总分']
    t_hit = sum(1 for k in template_keys if k in low)
    h_hit = sum(1 for k in homework_keys if k in low)
    if t_hit >= 2 and t_hit >= h_hit:
        return "教案/教学设计模板"
    if h_hit >= 2:
        return "学生作业/成绩/批改记录"
    if '.xlsx' in name.lower() or '.csv' in name.lower():
        return "数据表"
    if '.pdf' in name.lower():
        return "PDF文档"
    return "参考资料"


def _table_to_markdown(rows):
    rows = list(rows)
    if not rows:
        return ""
    rows = [[("" if v is None else str(v)).replace('|', '\\|').replace('\n', ' ') for v in r] for r in rows]
    max_cols = max(len(r) for r in rows)
    rows = [r + [''] * (max_cols - len(r)) for r in rows]
    header = rows[0]
    sep = ['---'] * max_cols
    body = rows[1:]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join(sep) + " |"]
    for r in body:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


@app.route("/api/upload", methods=["POST"])
def api_upload():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    if 'file' not in request.files:
        return jsonify({"error": "未收到上传文件字段 file"}), 400

    f = request.files['file']
    if not f or not f.filename:
        return jsonify({"error": "文件名为空"}), 400

    filename = os.path.basename(f.filename)
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXT:
        return jsonify({"error": f"不支持的格式 {ext}，支持：md / txt / docx / xlsx / csv / pdf"}), 400

    raw = f.read()
    if len(raw) > MAX_UPLOAD_BYTES:
        return jsonify({"error": f"文件过大 {len(raw)} bytes，上限 {MAX_UPLOAD_BYTES} bytes（20MB）"}), 400

    text = ""
    try:
        if ext in ('.txt', '.md'):
            text = _read_text_bytes(raw)
        elif ext == '.docx':
            bio = BytesIO(raw)
            doc = Document(bio)
            parts = []
            for p in doc.paragraphs:
                if p.text:
                    parts.append(p.text)
            for t in doc.tables:
                rows = [[c.text for c in row.cells] for row in t.rows]
                md = _table_to_markdown(rows)
                if md:
                    parts.append("\n[表格]\n" + md + "\n")
            text = "\n".join(parts)
        elif ext == '.csv':
            s = _read_text_bytes(raw)
            import csv as _csv
            reader = _csv.reader(s.splitlines())
            rows = list(reader)[:200]
            text = _table_to_markdown(rows) or (s[:MAX_TEXT_CHARS])
        elif ext == '.xlsx':
            import openpyxl
            bio = BytesIO(raw)
            wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
            sheets_text = []
            for ws in wb.worksheets:
                rows = []
                for idx, row in enumerate(ws.iter_rows(values_only=True)):
                    if idx > 60:
                        break
                    rows.append(list(row)[:16])
                md = _table_to_markdown(rows)
                if md:
                    sheets_text.append(f"# 工作表：{ws.title}\n{md}")
            text = "\n\n".join(sheets_text)
        elif ext == '.pdf':
            try:
                import PyPDF2
                bio = BytesIO(raw)
                reader = PyPDF2.PdfReader(bio)
                pages = []
                for i, page in enumerate(reader.pages[:30]):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    if t:
                        pages.append(f"\n---第 {i+1} 页---\n{t}")
                text = "\n".join(pages)
            except ImportError:
                return jsonify({"error": "PDF 需要先安装 PyPDF2：pip install PyPDF2"}), 400
    except Exception as ee:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"解析文件失败: {ee}"}), 500

    if not text or not text.strip():
        return jsonify({"error": "文件内容为空或无法识别文本"}), 400

    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n[原文过长已截断，仅保留前部分内容用于分析]"

    kind = _guess_file_kind(filename, text)
    preview = text[:600]

    return jsonify({
        "ok": True,
        "name": filename,
        "kind": kind,
        "chars": len(text),
        "preview": preview,
        "full_text": text,
    })


def _find_user_file(filename):
    user = get_current_user()
    if user:
        user_dir = get_user_output_dir(user['id'])
        fp = os.path.join(user_dir, filename)
        if os.path.isfile(fp):
            return fp
    if os.path.isdir(HISTORY_DIR):
        for sub in os.listdir(HISTORY_DIR):
            fp = os.path.join(HISTORY_DIR, sub, filename)
            if os.path.isfile(fp):
                return fp
    fp = os.path.join(OUTPUT_DIR, filename)
    if os.path.isfile(fp):
        return fp
    return None


@app.route("/api/download/<path:filename>")
def download(filename):
    from urllib.parse import unquote, quote as _qd
    filename = unquote(filename)
    filepath = _find_user_file(filename)
    if not filepath:
        return jsonify({"error": "文件不存在"}), 404

    ext = os.path.splitext(filename)[1].lower()
    mimetype = MIME_MAP.get(ext, 'application/octet-stream')
    try:
        safe_ascii = filename.encode('ascii', errors='ignore').decode('ascii') or ('download' + ext)
    except Exception:
        safe_ascii = 'download' + ext
    try:
        resp = send_file(filepath, as_attachment=True, download_name=safe_ascii, mimetype=mimetype)
        resp.headers['X-Filename-Encoded'] = _qd(filename)
        resp.headers['X-Filename'] = safe_ascii
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as ee:
        print(f"[download] 发送文件异常: {ee}")
        import traceback; traceback.print_exc()
        return jsonify({"error": f"文件下载失败: {ee}"}), 500


@app.route("/api/export/word", methods=["POST"])
def export_to_word():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = request.json
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    markdown_content = data.get("content", "").strip()
    title = str(data.get("title", "教案")).strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])

    try:
        word_buffer = markdown_to_word(markdown_content, title)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:20] or "untitled"
        filename = f"教案_{safe_title}_{timestamp}.docx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(word_buffer.getvalue())

        from urllib.parse import quote as _q
        from unicodedata import normalize as _norm
        try:
            _ascii_fn = filename.encode('ascii', errors='ignore').decode('ascii') or 'document.docx'
        except Exception:
            _ascii_fn = 'document.docx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn,
                         mimetype=MIME_MAP['.docx'])
        resp.headers['X-Filename-Encoded'] = _q(filename)
        resp.headers['X-Filename'] = _ascii_fn
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"导出失败: {str(e)}"}), 500

@app.route("/api/export/ppt", methods=["POST"])
def export_to_ppt():
    user = get_current_user()
    if not user:
        return jsonify({"error": "请先登录"}), 401

    data = request.json
    if not data or not isinstance(data, dict):
        return jsonify({"error": "Invalid request body"}), 400

    markdown_content = data.get("content", "").strip()
    title = str(data.get("title", "课件")).strip()

    if not markdown_content:
        return jsonify({"error": "内容不能为空"}), 400

    if len(markdown_content) > 500000:
        return jsonify({"error": "内容过长，无法导出"}), 400

    user_output_dir = get_user_output_dir(user['id'])

    try:
        ppt_buffer = markdown_to_ppt(markdown_content, title)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[<>:"/\\|?*]', '', title)[:20] or "untitled"
        filename = f"课件_{safe_title}_{timestamp}.pptx"
        filepath = os.path.join(user_output_dir, filename)

        with open(filepath, "wb") as f:
            f.write(ppt_buffer.getvalue())

        from urllib.parse import quote as _q2
        try:
            _ascii_fn2 = filename.encode('ascii', errors='ignore').decode('ascii') or 'slides.pptx'
        except Exception:
            _ascii_fn2 = 'slides.pptx'
        resp = send_file(filepath, as_attachment=True, download_name=_ascii_fn2,
                         mimetype=MIME_MAP['.pptx'])
        resp.headers['X-Filename-Encoded'] = _q2(filename)
        resp.headers['X-Filename'] = _ascii_fn2
        resp.headers['Access-Control-Expose-Headers'] = 'X-Filename, X-Filename-Encoded, Content-Disposition'
        return resp
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({"error": f"导出 PPT 失败: {str(e)}"}), 500

MODULES_DATA = {
    "ppt_generate": {
        "name": "PPT生成",
        "emoji": "📊",
        "short_desc": "一键生成教学课件PPT",
        "description": "AI智能生成完整教学课件，支持从课程主题直接输出结构化PPT大纲，并可一键导出为 .pptx 文件。自动包含封面、教学目标、知识点讲解、案例分析、课堂小结、课后作业等标准课件结构。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "生成一份《荷塘月色》的语文课件PPT，包含教学目标、作者介绍、课文赏析、重点语句分析、课堂练习、作业布置，共15-20页",
        "features": [
            "� 标准课件结构（封面/目标/新知/小结/作业）",
            "📝 每页幻灯片内容完整撰写（标题+要点）",
            "🎨 Markdown格式规范排版，方便二次编辑",
            "⬇️ 一键导出 .pptx 文件，可直接打开使用",
            "� 自动适配学科特点（文科赏析/理科例题/工科原理）",
            "⏱ 支持指定页数与各环节时间分配"
        ],
        "examples": [
            {"title": "语文课件：荷塘月色", "desc": "朱自清散文，含作者背景+课文赏析+修辞分析", "prompt": "生成一份《荷塘月色》的语文课件PPT，适合高中，包含教学目标、作者介绍、课文逐段赏析、修辞手法分析、课堂练习、作业布置，共15-20页"},
            {"title": "数学课件：二次函数", "desc": "初中数学，含图像性质+典型例题+课堂练习", "prompt": "生成初中数学《二次函数的图像与性质》课件PPT，包含概念讲解、5个典型例题解析、课堂练习、课后作业，约18页"},
            {"title": "历史课件：辛亥革命", "desc": "高中历史，含背景脉络+重大事件+意义启示", "prompt": "生成高中历史《辛亥革命》课件PPT，包含历史背景、武昌起义经过、中华民国成立、历史意义与局限、思考题，约20页"}
        ],
        "knowledge_points": ["课件结构", "教学目标", "知识点导入", "案例分析", "课堂互动", "例题解析", "归纳小结", "作业布置", "板书设计", "时间分配"]
    },
    "ppt_outline": {
        "name": "PPT讲解纲要",
        "emoji": "�",
        "short_desc": "每页PPT的逐页讲解稿",
        "description": "根据PPT大纲为每一页幻灯片生成详细的教师讲解词。包含开场白设计、知识点过渡衔接语、重难点突出提示、提问互动设计、随堂练习引导语等，帮助教师脱稿流畅完成整堂课。",
        "subject": "通用",
        "grade_levels": ["小学", "初中", "高中", "职业院校", "大学"],
        "default_prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、互动提示、过渡语句、重点强调部分",
        "features": [
            "🎤 逐页讲解词撰写，教师可直接照着读",
            "🔗 页面之间设计自然衔接过渡语句",
            "❓ 穿插课堂提问与互动设计（教师问+预期答）",
            "⭐ 标注重点知识点的语气语调提示（重读、停顿）",
            "⏰ 每页推荐讲解时长，合理分配45分钟",
            "🧑‍🏫 模拟真实课堂的教态与现场感"
        ],
        "examples": [
            {"title": "语文：荷塘月色讲解稿", "desc": "逐页教师用语+情感渲染引导", "prompt": "为高中语文《荷塘月色》PPT课件生成逐页讲解纲要，包含每页幻灯片的教师讲解词、课堂提问设计、过渡语句、重点部分标注语气要求"},
            {"title": "数学：二次函数讲解稿", "desc": "概念引入+例题推导步骤+易错提醒", "prompt": "生成初中数学《二次函数图像与性质》每一页PPT的教师讲解纲要，包含概念引入话术、例题逐步推导讲解、学生易错点提醒、课堂练习的处理方式"},
            {"title": "班会：防溺水安全教育稿", "desc": "案例警示+互动问答+总结升华", "prompt": "生成中小学防溺水主题班会PPT的逐页讲解纲要，包含真实案例描述、互动提问、自救技巧讲解、最后的倡议环节，语言符合中小学生认知"}
        ],
        "knowledge_points": ["开场白设计", "过渡衔接", "重难点提示", "互动提问", "案例讲述", "语气语调", "时间把控", "课堂节奏", "情感升华", "结尾总结"]
    },
    "lesson_plan": {
        "name": "教案生成",
        "emoji": "�",
        "short_desc": "标准格式45分钟完整教案",
        "description": "按中小学教案标准格式自动生成完整45分钟教学设计。包含教学目标（三维/核心素养）、教学重难点、教学方法、教学过程（导入/新授/练习/小结/作业）、板书设计、教学反思模板等全部要素。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "中职", "高职"],
        "default_prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程、板书设计、作业布置、教学反思",
        "features": [
            "🎯 三维教学目标（知识与技能/过程与方法/情感态度）",
            "🎯 核心素养目标适配新课标要求",
            "❗ 教学重点与教学难点明确区分",
            "� 详细教学过程：5-6个环节精确到分钟",
            "🎨 板书设计框架（主板书+副板书）",
            "� 配套作业设计（基础+拓展分层）",
            "� 预留教学反思栏（课堂生成问题记录）"
        ],
        "examples": [
            {"title": "语文：背影教案", "desc": "朱自清经典·情感教育·细节描写", "prompt": "生成一份初中语文《背影》的完整45分钟教案，包含三维教学目标、教学重难点、详细教学过程（各环节含时间分配）、板书设计、分层作业布置、教学反思栏"},
            {"title": "物理：浮力教案", "desc": "实验探究·阿基米德原理·分层练习", "prompt": "生成初中物理《阿基米德原理》45分钟教案，包含实验探究环节设计、演示实验步骤、学生分组活动安排、分层练习设计、板书设计"},
            {"title": "英语：Travel Plan教案", "desc": "听说课型·情境交际·任务型教学", "prompt": "生成初中英语听说课《Travel Plans》完整教案，包含Warm-up、Pre-listening、While-listening、Post-speaking、Summary & Homework环节，配师生对话示例"}
        ],
        "knowledge_points": ["三维目标", "核心素养", "教学重难点", "学情分析", "教学方法", "教学过程", "时间分配", "板书设计", "分层作业", "教学反思"]
    },
    "exercises": {
        "name": "习题生成",
        "emoji": "📝",
        "short_desc": "分层习题+详细答案解析",
        "description": "根据知识点自动生成多难度层次练习题，支持选择题、填空题、判断题、解答题、应用题等多种题型。每道题目均附标准答案与详细解题步骤，可用于课堂练习、课后作业、单元测验。",
        "subject": "全学科",
        "grade_levels": ["小学低年级", "小学高年级", "初中", "高中", "中职", "高职"],
        "default_prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用，所有题目附答案和详细解析",
        "features": [
            "� 三级难度分层：基础巩固·能力提高·拓展探究",
            "✅ 题型丰富：选择/填空/判断/解答/实验探究/应用题",
            "� 每题附标准答案 + 详细解题步骤",
            "� 标注每题考查知识点对应章节",
            "⏱ 推荐完成时长与分值设置",
            "� 可一键组装为标准试卷（卷头+姓名栏）"
        ],
        "examples": [
            {"title": "数学：勾股定理20题", "desc": "8基础+8提高+4拓展，附详解", "prompt": "为初中数学《勾股定理》生成20道分层练习题：基础题8道选择+填空，提高题8道计算+证明，拓展题4道综合应用题，所有题目附答案和详细解析步骤"},
            {"title": "英语：一般过去时语法", "desc": "30道选择+填空+句型转换，解析全", "prompt": "生成初中英语语法《一般过去时》专项练习题30道：单项选择15道，用所给动词适当形式填空10道，句型转换5道，附答案和详细解析"},
            {"title": "化学：酸碱盐专题", "desc": "选择+填空+推断+实验探究四大题型", "prompt": "生成初中化学《酸碱盐》单元练习题，包含选择题10道、填空题6道、物质推断题2道、实验探究题2道，附答案与解析"}
        ],
        "knowledge_points": ["基础巩固", "能力提高", "拓展探究", "选择题", "填空题", "解答题", "实验探究", "应用题", "答案解析", "考点分布"]
    },
    "homework_analysis": {
        "name": "作业分析",
        "emoji": "�",
        "short_desc": "作业批改与学情数据分析",
        "description": "上传作业统计数据或描述作业情况，AI自动生成完整的作业分析报告。包含：正确率统计、高频错题归因、典型错解展示、知识点掌握情况雷达图、学困生针对性辅导建议、下节课教学调整策略等。",
        "subject": "全学科",
        "grade_levels": ["小学", "初中", "高中", "职业院校"],
        "default_prompt": "请为一次高中数学《函数单调性》课后作业生成作业分析报告：全班45人，平均得分率68%，错误集中在：含参数的单调性讨论、复合函数单调性、实际应用最值问题，请给出详细分析与辅导建议",
        "features": [
            "� 整体成绩统计：平均分/得分率/分数段分布",
            "❌ 高频错题 TOP 5 归因分析（概念/审题/计算/方法）",
            "� 典型错解案例展示与正解对比",
            "🎯 知识点掌握度雷达图（熟练/一般/薄弱）",
            "👨‍🎓 分层辅导建议：学困生、中等生、优等生",
            "✏️ 下一步教学调整：需补讲的知识点、课堂策略",
            "📑 讲评课时分配建议与讲评例题推荐"
        ],
        "examples": [
            {"title": "数学函数单调性作业分析", "desc": "得分率68%·参数讨论错误高发", "prompt": "为高中数学《函数单调性》课后作业生成详细作业分析报告：全班45人，平均得分率68%，主要错误：①含参数单调性分类讨论（失分率52%）②复合函数同增异减应用（失分率41%）③实际问题求最值忽略定义域（失分率37%），请给出错因分析、辅导建议、教学调整策略"},
            {"title": "英语完形填空作业分析", "desc": "上下文逻辑·固定搭配是薄弱点", "prompt": "分析一次八年级英语完形填空作业情况：全班40人，平均正确率58%；错题中上下文逻辑理解占45%，固定搭配与词组占30%，词汇辨析占15%，语法占10%。请生成完整作业讲评方案"},
            {"title": "物理电路作业错误分析", "desc": "串并联识别+欧姆定律综合应用薄弱", "prompt": "为初三物理《欧姆定律在串并联电路中的应用》作业做分析：全班48人，平均得分62分；常见错：复杂电路等效化简错误、电表测量对象判断错误、比例计算错误。生成作业分析与讲评建议"}
        ],
        "knowledge_points": ["得分率统计", "错题归因", "典型错解", "知识薄弱点", "分层辅导", "学困生帮扶", "教学调整", "讲评课时", "补充练习", "家校沟通"]
    },
    "history": {
        "name": "历史记录",
        "emoji": "🕒",
        "short_desc": "查看所有AI生成记录与历史会话",
        "description": "查看当前账号下所有历史生成记录：历史会话聊天、PPT课件、教案、习题、作业分析等全部生成内容。支持按时间、按模块筛选，快速找到过往生成内容，可再次打开、复用、编辑或重新生成。",
        "subject": "系统功能",
        "grade_levels": ["全部学段"],
        "default_prompt": "（点击下方会话卡片直接查看历史记录）",
        "features": [
            "🕒 按时间倒序展示所有历史会话",
            "🔍 支持按模块类型筛选（PPT/教案/习题/作业等）",
            "� 可快速查看每条记录的生成时间、标题、内容摘要",
            "♻️ 一键复用：把历史内容重新发送，基于结果再修改",
            "🗂 可导出历史记录为 Word / 文本文件",
            "� 查看个人使用统计（生成次数、常用模块等）"
        ],
        "examples": [
            {"title": "查看最近7天生成记录", "desc": "快速找到上周生成的PPT并修改", "prompt": "__HISTORY__VIEW__"},
            {"title": "按模块筛选：只看教案", "desc": "查找所有历史生成的教案", "prompt": "__HISTORY_VIEW__ lesson_plan"},
            {"title": "统计本月使用情况", "desc": "查看自己备课模块使用分布", "prompt": "__HISTORY_STATS__"}
        ],
        "knowledge_points": ["聊天记录", "生成文件", "按时间筛选", "按模块筛选", "内容搜索", "记录复用", "批量导出", "使用统计", "收藏夹", "回收站"]
    }
}

@app.route("/module/<module_id>")
def module_detail(module_id):
    user = get_current_user()
    if not user:
        return redirect(url_for('login'))
    if module_id not in MODULES_DATA:
        return redirect(url_for('index'))
    module = MODULES_DATA[module_id]
    all_modules = [{"id": k, "name": v["name"], "emoji": v["emoji"]} for k, v in MODULES_DATA.items()]
    history_sessions = []
    if module_id == "history":
        try:
            MODULE_META = {
                "ppt_generate":         {"name": "PPT生成",     "emoji": "📊", "color": "#667eea", "gradient": "linear-gradient(135deg,#667eea 0%,#764ba2 100%)"},
                "ppt_outline":          {"name": "PPT讲解纲要", "emoji": "📝", "color": "#f093fb", "gradient": "linear-gradient(135deg,#f093fb 0%,#f5576c 100%)"},
                "lesson_plan_generate": {"name": "教案生成",     "emoji": "📘", "color": "#4facfe", "gradient": "linear-gradient(135deg,#4facfe 0%,#00f2fe 100%)"},
                "exercise_generate":    {"name": "习题生成",     "emoji": "📚", "color": "#43e97b", "gradient": "linear-gradient(135deg,#43e97b 0%,#38f9d7 100%)"},
                "homework_analysis":    {"name": "作业分析",     "emoji": "📊", "color": "#fa709a", "gradient": "linear-gradient(135deg,#fa709a 0%,#fee140 100%)"},
                "history":              {"name": "历史记录",     "emoji": "📜", "color": "#a8edea", "gradient": "linear-gradient(135deg,#a8edea 0%,#fed6e3 100%)"},
                "general":              {"name": "自由聊天",     "emoji": "💬", "color": "#c3cfe2", "gradient": "linear-gradient(135deg,#c3cfe2 0%,#f5f7fa 100%)"},
            }
            stats = database.get_user_sessions_with_stats(user['id'])
            for s in stats:
                mid = s.get("module_id") or "general"
                if mid not in MODULE_META:
                    mid = "general"
                meta = MODULE_META[mid]
                # 汇总文件徽章：按 kind 去重后给前端
                files = s.get("files") or []
                badge_summary = {}
                for f in files:
                    k = f.get("kind") or "文件"
                    badge_summary.setdefault(k, {"count": 0, "samples": []})
                    if badge_summary[k]["count"] < 3:
                        badge_summary[k]["samples"].append(f.get("name") or "")
                    badge_summary[k]["count"] += 1
                badges = []
                for k, v in badge_summary.items():
                    ext_icon = {"PPT": "📊", "教案": "📘", "练习册": "📚", "作业分析": "📈", "报告": "📄", "文件": "📄"}.get(k, "📄")
                    badges.append({
                        "kind": k,
                        "icon": ext_icon,
                        "count": v["count"],
                        "sample": (v["samples"][0] or "")[:50]
                    })
                history_sessions.append({
                    "id": s["session_id"],
                    "pk": s.get("pk_id") or "",
                    "title": s.get("title") or "(空对话)",
                    "created_at": s.get("created_at") or "",
                    "updated_at": s.get("updated_at") or "",
                    "time": s.get("updated_at") or s.get("created_at") or "",
                    "msg_count": s.get("total_count") or 0,
                    "user_count": s.get("user_count") or 0,
                    "assistant_count": s.get("assistant_count") or 0,
                    "file_count": s.get("file_count") or 0,
                    "module_id": mid,
                    "module_name": meta["name"],
                    "module_emoji": meta["emoji"],
                    "module_color": meta["color"],
                    "module_gradient": meta["gradient"],
                    "preview": s.get("user_preview") or s.get("ai_preview") or "",
                    "user_preview": s.get("user_preview") or "",
                    "ai_preview": s.get("ai_preview") or "",
                    "badges": badges,
                    "files": files[:8],
                })
            # 按最后更新时间倒序（database 已经排过，这里再保证一次）
            history_sessions.sort(key=lambda x: x.get("time") or "", reverse=True)
        except Exception as e:
            import traceback; traceback.print_exc()
            print("[history] 读取会话失败：", e)
    return render_template(
        "module_detail.html",
        module=module,
        module_id=module_id,
        model=MODEL,
        username=user['username'],
        all_modules=all_modules,
        history_sessions=history_sessions
    )

if __name__ == "__main__":
    print("""
============================================
       BeiKe Assistant (Web) - 登录版 2.0
       http://localhost:5002
       Ctrl+C to exit
============================================
""")
    import webbrowser
    threading.Timer(1.5, lambda: webbrowser.open("http://127.0.0.1:5002/")).start()
    app.run(host="0.0.0.0", port=5002, debug=False, threaded=True)
